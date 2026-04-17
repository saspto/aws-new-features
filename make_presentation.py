"""
Build the AWS Features Tracker innovation demo PowerPoint presentation.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import copy
from lxml import etree

ASSETS = "/workshop/aws-new-features/slides_assets"
OUT    = "/workshop/aws-new-features/AWS_Features_Tracker_Innovation_Demo.pptx"

# ── brand colours ────────────────────────────────────────────────────────────
AWS_ORANGE  = RGBColor(0xFF, 0x99, 0x00)   # #FF9900
AWS_DARK    = RGBColor(0x23, 0x2F, 0x3E)   # #232F3E
AWS_BLUE    = RGBColor(0x14, 0x6E, 0xB4)   # #146EB4
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF8, 0xFA, 0xFC)
MID_GRAY    = RGBColor(0x94, 0xA3, 0xB8)
DARK_GRAY   = RGBColor(0x1E, 0x29, 0x3B)
ACCENT_TEAL = RGBColor(0x0F, 0xB5, 0xD4)
GREEN       = RGBColor(0x16, 0xA3, 0x4A)

# ── slide dimensions (16:9) ──────────────────────────────────────────────────
SW = Inches(13.33)
SH = Inches(7.5)

prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH

BLANK = prs.slide_layouts[6]   # completely blank

# ─────────────────────────────────────────────────────────────────────────────
# Helper utilities
# ─────────────────────────────────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width_pt=0):
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, l, t, w, h, text, font_size, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return tb

def add_paragraph(tf, text, font_size, bold=False, color=DARK_GRAY,
                  align=PP_ALIGN.LEFT, space_before=0, italic=False):
    from pptx.util import Pt as _Pt
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = _Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size  = _Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return p

def add_icon_bullet(slide, l, t, w, h, icon_char, text, font_size=12,
                    icon_color=AWS_ORANGE, text_color=DARK_GRAY):
    """Dot + text row helper."""
    # dot
    add_rect(slide, l, t + Pt(font_size*0.3), Inches(0.12), Inches(0.12), fill_color=icon_color)
    tb = slide.shapes.add_textbox(l + Inches(0.22), t, w - Inches(0.22), h)
    tf = tb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = text_color
    return tb

def add_pill(slide, l, t, w, h, text, bg, fg, font_size=9):
    add_rect(slide, l, t, w, h, fill_color=bg)
    add_textbox(slide, l, t + Pt(1), w, h, text, font_size, color=fg, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# Slide helpers
# ─────────────────────────────────────────────────────────────────────────────
def title_accent_bar(slide):
    add_rect(slide, 0, SH - Inches(0.12), SW, Inches(0.12), fill_color=AWS_ORANGE)

def section_label(slide, text, l=Inches(0.55), t=Inches(0.28)):
    add_textbox(slide, l, t, Inches(4), Inches(0.3),
                text, 9, bold=True, color=AWS_ORANGE, align=PP_ALIGN.LEFT)

def divider(slide, t, color=RGBColor(0xE2,0xE8,0xF0)):
    add_rect(slide, Inches(0.55), t, SW - Inches(1.1), Pt(1), fill_color=color)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title / Hero
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)

# dark background
add_rect(slide, 0, 0, SW, SH, fill_color=AWS_DARK)

# orange accent strip left
add_rect(slide, 0, 0, Inches(0.18), SH, fill_color=AWS_ORANGE)

# subtle grid lines (decorative)
for i in range(1, 9):
    add_rect(slide, Inches(i*13.33/8), 0, Pt(0.5), SH,
             fill_color=RGBColor(0x2D, 0x3F, 0x55))

# AWS logo area — orange square + text
add_rect(slide, Inches(0.65), Inches(1.1), Inches(0.72), Inches(0.48), fill_color=AWS_ORANGE)
add_textbox(slide, Inches(0.68), Inches(1.12), Inches(0.66), Inches(0.44),
            "aws", 22, bold=True, color=AWS_DARK, align=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0.65), Inches(1.72), Inches(10), Inches(0.55),
            "INNOVATION DEMO  ·  AWS ADMINISTRATION TEAM", 11,
            color=AWS_ORANGE, bold=True)

# Main headline
add_textbox(slide, Inches(0.65), Inches(2.4), Inches(10), Inches(1.05),
            "AWS Features Tracker", 52, bold=True, color=WHITE)

# Sub-headline
add_textbox(slide, Inches(0.65), Inches(3.52), Inches(9.5), Inches(0.55),
            "AI-Powered Discovery, Summarization & Search of AWS Announcements", 22,
            color=RGBColor(0xCB, 0xD5, 0xE1), italic=True)

# Three value pills
pill_labels = ["Real-time RSS ingestion", "Claude Haiku AI summaries", "Searchable web app"]
for i, lbl in enumerate(pill_labels):
    px = Inches(0.65 + i * 3.3)
    add_rect(slide, px, Inches(4.3), Inches(3.0), Inches(0.4),
             fill_color=RGBColor(0x14, 0x6E, 0xB4))
    add_textbox(slide, px, Inches(4.31), Inches(3.0), Inches(0.38),
                lbl, 12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Bottom bar
add_rect(slide, 0, SH - Inches(0.5), SW, Inches(0.5), fill_color=RGBColor(0x14,0x21,0x30))
add_textbox(slide, 0, SH - Inches(0.46), SW, Inches(0.42),
            "Presented by the Cloud Administration Team  ·  April 2026", 10,
            color=MID_GRAY, align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem / Opportunity
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, SW, SH, fill_color=LIGHT_GRAY)
add_rect(slide, 0, 0, SW, Inches(1.1), fill_color=AWS_DARK)
title_accent_bar(slide)

add_textbox(slide, Inches(0.55), Inches(0.22), Inches(10), Inches(0.56),
            "The Challenge: Staying Ahead of AWS Innovation", 28, bold=True, color=WHITE)
add_textbox(slide, Inches(0.55), Inches(0.75), Inches(10), Inches(0.28),
            "AWS releases 2,000+ new features per year — how do admins keep up?", 13,
            color=AWS_ORANGE, italic=True)

# Stats row
stats = [
    ("2,000+", "new AWS\nfeatures/year"),
    ("~5 min", "avg time to\nread each post"),
    ("166 hrs", "of reading\nannually"),
    ("0%", "of that time\nbillable"),
]
for i, (num, lbl) in enumerate(stats):
    bx = Inches(0.55 + i * 3.12)
    add_rect(slide, bx, Inches(1.25), Inches(2.9), Inches(1.5), fill_color=WHITE,
             line_color=RGBColor(0xE2,0xE8,0xF0), line_width_pt=1)
    add_textbox(slide, bx, Inches(1.35), Inches(2.9), Inches(0.72),
                num, 36, bold=True, color=AWS_ORANGE, align=PP_ALIGN.CENTER)
    add_textbox(slide, bx, Inches(2.05), Inches(2.9), Inches(0.58),
                lbl, 12, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# Pain points
divider(slide, Inches(2.95))
add_textbox(slide, Inches(0.55), Inches(3.05), Inches(6), Inches(0.32),
            "Pain Points", 15, bold=True, color=DARK_GRAY)

pains = [
    "Manual RSS monitoring across dozens of categories — easy to miss critical announcements",
    "No searchable history — can't quickly answer 'did AWS update Lambda timeouts this year?'",
    "Raw HTML descriptions are noisy — critical details buried in marketing language",
    "No way to filter by relevance to our specific services (EC2, RDS, Lambda, EKS, S3…)",
    "Team knowledge is siloed — each admin tracks different services independently",
]
for i, pain in enumerate(pains):
    add_icon_bullet(slide, Inches(0.55), Inches(3.48 + i*0.52), Inches(8), Inches(0.45),
                    "•", pain, 12, icon_color=RGBColor(0xEF,0x44,0x44), text_color=DARK_GRAY)

# Opportunity box
add_rect(slide, Inches(9.05), Inches(3.05), Inches(3.72), Inches(3.45),
         fill_color=RGBColor(0xFF, 0xF7, 0xED), line_color=AWS_ORANGE, line_width_pt=2)
add_textbox(slide, Inches(9.15), Inches(3.15), Inches(3.52), Inches(0.36),
            "The Opportunity", 14, bold=True, color=AWS_ORANGE)
add_textbox(slide, Inches(9.15), Inches(3.55), Inches(3.52), Inches(2.7),
            "Automate the discovery, summarization, and cataloguing of every AWS announcement "
            "— so admins spend minutes on insights instead of hours on reading.",
            12, color=DARK_GRAY, wrap=True)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Solution Overview / Architecture
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, SW, SH, fill_color=AWS_DARK)
title_accent_bar(slide)

section_label(slide, "SOLUTION ARCHITECTURE")
add_textbox(slide, Inches(0.55), Inches(0.5), Inches(10), Inches(0.56),
            "AWS Features Tracker — How It Works", 28, bold=True, color=WHITE)

# Architecture flow boxes
FLOW_Y  = Inches(1.35)
BOX_H   = Inches(1.1)
ARROW_Y = FLOW_Y + BOX_H/2

steps = [
    ("EventBridge\nScheduler", "Every 6 hours", AWS_ORANGE),
    ("Fetcher\nLambda", "Parses RSS XML\nDeduplicates GUIDs", AWS_BLUE),
    ("SQS Queue", "Reliable delivery\nDLQ on failure", ACCENT_TEAL),
    ("Summarizer\nLambda", "Claude Haiku\nvia Bedrock", GREEN),
    ("DynamoDB\n+ S3", "Persistent store\nJSON archive", RGBColor(0x8B,0x5C,0xF6)),
]
BOX_W = Inches(2.18)
GAP   = Inches(0.12)
START = Inches(0.35)

for i, (name, desc, color) in enumerate(steps):
    bx = START + i * (BOX_W + GAP + Inches(0.3))
    add_rect(slide, bx, FLOW_Y, BOX_W, BOX_H, fill_color=color)
    add_textbox(slide, bx, FLOW_Y + Inches(0.1), BOX_W, Inches(0.44),
                name, 13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, bx, FLOW_Y + Inches(0.54), BOX_W, Inches(0.5),
                desc, 10, color=RGBColor(0xE2,0xE8,0xF0), align=PP_ALIGN.CENTER)
    if i < len(steps)-1:
        ax = bx + BOX_W + Inches(0.05)
        add_textbox(slide, ax, ARROW_Y - Inches(0.12), Inches(0.3), Inches(0.28),
                    "→", 18, bold=True, color=AWS_ORANGE, align=PP_ALIGN.CENTER)

# Second row: API layer
add_textbox(slide, Inches(0.35), Inches(2.7), Inches(12.5), Inches(0.3),
            "▼  API & Frontend Layer", 11, bold=True, color=MID_GRAY)

API_Y = Inches(3.1)
api_steps = [
    ("API Gateway\n(HTTP API)", "REST endpoints\nNo auth — read-only", AWS_BLUE),
    ("API Lambda", "/features /services\n/types /health", RGBColor(0x0F,0xB5,0xD4)),
    ("CloudFront CDN", "Global edge caching\nSPA routing", AWS_ORANGE),
    ("React SPA\n(S3)", "TypeScript + Vite\nTailwind CSS", RGBColor(0x8B,0x5C,0xF6)),
]
for i, (name, desc, color) in enumerate(api_steps):
    bx = START + i * (BOX_W + GAP + Inches(0.3))
    add_rect(slide, bx, API_Y, BOX_W, BOX_H, fill_color=color)
    add_textbox(slide, bx, API_Y + Inches(0.1), BOX_W, Inches(0.44),
                name, 13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, bx, API_Y + Inches(0.54), BOX_W, Inches(0.5),
                desc, 10, color=RGBColor(0xE2,0xE8,0xF0), align=PP_ALIGN.CENTER)
    if i < len(api_steps)-1:
        ax = bx + BOX_W + Inches(0.05)
        add_textbox(slide, ax, API_Y + BOX_H/2 - Inches(0.12), Inches(0.3), Inches(0.28),
                    "→", 18, bold=True, color=AWS_ORANGE, align=PP_ALIGN.CENTER)

# Key tech row
add_textbox(slide, Inches(0.35), Inches(4.45), Inches(12.5), Inches(0.3),
            "Tech Stack", 11, bold=True, color=MID_GRAY)
tech = ["Python 3.12 · Lambda arm64",
        "Amazon Bedrock · Claude Haiku",
        "DynamoDB on-demand",
        "React 18 + TypeScript + Vite",
        "AWS SAM (IaC)",
        "~$0.31/month"]
for i, t in enumerate(tech):
    tx = Inches(0.35 + i * 2.17)
    add_rect(slide, tx, Inches(4.8), Inches(2.05), Inches(0.4),
             fill_color=RGBColor(0x1E,0x3A,0x52))
    add_textbox(slide, tx, Inches(4.81), Inches(2.05), Inches(0.38),
                t, 9, color=RGBColor(0xCB,0xD5,0xE1), align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Key Features for AWS Admins
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, SW, SH, fill_color=LIGHT_GRAY)
add_rect(slide, 0, 0, SW, Inches(1.1), fill_color=AWS_DARK)
title_accent_bar(slide)

add_textbox(slide, Inches(0.55), Inches(0.2), Inches(11), Inches(0.56),
            "Built for AWS Admins Who Research & Analyze", 27, bold=True, color=WHITE)
add_textbox(slide, Inches(0.55), Inches(0.76), Inches(11), Inches(0.28),
            "Every feature is designed around the admin workflow — discover, evaluate, act.", 13,
            color=AWS_ORANGE, italic=True)

features = [
    (AWS_ORANGE, "Automated Discovery",
     "Fetches https://aws.amazon.com/new/feed/ every 6 hours via EventBridge Scheduler. "
     "Zero manual effort — new announcements appear within hours of posting.",
     ["No manual RSS monitoring ever again",
      "Deduplication ensures no double-processing",
      "Stores raw XML timestamp in SSM Parameter Store"]),
    (AWS_BLUE, "AI-Generated Summaries",
     "Claude Haiku (Amazon Bedrock) summarizes each announcement into 3-5 actionable sentences "
     "focused on what changed, who it benefits, and key technical limits.",
     ["Feature type auto-classified: new-feature, GA, preview, deprecation...",
      "3 key bullet points extracted per announcement",
      "Async processing via SQS — never blocks ingestion"]),
    (GREEN, "Powerful Search & Filters",
     "Full-text search across title, summary, and service. Filter by AWS service, feature type, "
     "and date range. All state is stored in the URL — results are bookmarkable.",
     ["Search 'lambda timeout' and find it in seconds",
      "Filter by service: amazon-ec2, aws-lambda, amazon-rds…",
      "Date range: 'what changed in Q1 2026?'"]),
    (ACCENT_TEAL, "Feature Type Classification",
     "Bedrock classifies every announcement into one of 8 types so admins can instantly filter "
     "for what matters most to their team.",
     ["new-feature · enhancement · preview · general-availability",
      "deprecation · price-change · region-expansion · integration",
      "Color-coded badges for instant visual scanning"]),
    (RGBColor(0x8B,0x5C,0xF6), "Shareable URLs & Detail Pages",
     "Every search query and feature detail has a unique URL. Share a filtered view with your "
     "team or bookmark a specific announcement for later review.",
     ["URL encodes all filters — copy/paste to share with colleagues",
      "Detail page shows full summary, all key points, and AWS link",
      "Back button preserves search context"]),
    (RGBColor(0xEF,0x44,0x44), "Operational Reliability",
     "SQS DLQ catches failed summaries, CloudWatch alarms fire on errors, "
     "and IAM least-privilege policies protect all resources.",
     ["DLQ retries up to 3× before dead-lettering",
      "CloudWatch alarms on Summarizer errors and DLQ depth",
      "All data retained indefinitely in DynamoDB + S3"]),
]

COL_W = Inches(4.1)
ROW_H = Inches(1.8)
for i, (color, title, body, bullets) in enumerate(features):
    col = i % 3
    row = i // 3
    bx = Inches(0.45 + col * (COL_W + Inches(0.12)))
    by = Inches(1.3 + row * (ROW_H + Inches(0.1)))

    # card
    add_rect(slide, bx, by, COL_W, ROW_H, fill_color=WHITE,
             line_color=RGBColor(0xE2,0xE8,0xF0), line_width_pt=1)
    # color bar
    add_rect(slide, bx, by, Inches(0.08), ROW_H, fill_color=color)

    add_textbox(slide, bx+Inches(0.15), by+Inches(0.1), COL_W-Inches(0.22), Inches(0.3),
                title, 13, bold=True, color=DARK_GRAY)
    add_textbox(slide, bx+Inches(0.15), by+Inches(0.42), COL_W-Inches(0.22), Inches(0.56),
                body, 9.5, color=RGBColor(0x64,0x74,0x8B), wrap=True)
    for j, bul in enumerate(bullets[:2]):
        add_icon_bullet(slide, bx+Inches(0.15), by+Inches(1.02+j*0.32),
                        COL_W-Inches(0.22), Inches(0.28),
                        "•", bul, 9, icon_color=color,
                        text_color=DARK_GRAY)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Screenshot: Main Search Page
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, SW, SH, fill_color=AWS_DARK)
title_accent_bar(slide)

section_label(slide, "LIVE DEMO — SEARCH & DISCOVERY")
add_textbox(slide, Inches(0.55), Inches(0.44), Inches(9), Inches(0.48),
            "Main Dashboard — Browse & Search All AWS Announcements", 22, bold=True, color=WHITE)

# screenshot
slide.shapes.add_picture(f"{ASSETS}/screen_search.png",
                          Inches(0.35), Inches(1.0), Inches(12.65), Inches(5.95))

# callout labels overlaid
callouts = [
    (Inches(0.4),  Inches(1.05), "Live search bar with 300ms debounce"),
    (Inches(0.4),  Inches(1.55), "Smart filters: service, type, date range, sort order"),
    (Inches(0.4),  Inches(2.1),  "Result count always visible"),
    (Inches(9.1),  Inches(1.05), "Color-coded feature type badges"),
    (Inches(9.1),  Inches(1.55), "AI-extracted key points on hover"),
    (Inches(9.1),  Inches(2.1),  "Direct AWS announcement link"),
]
for cx, cy, text in callouts:
    add_rect(slide, cx, cy, Inches(3.1), Inches(0.3), fill_color=RGBColor(0x23,0x2F,0x3E))
    add_textbox(slide, cx+Inches(0.08), cy+Inches(0.04), Inches(3.0), Inches(0.26),
                "▶  " + text, 9, color=AWS_ORANGE)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Screenshot: Feature Detail Page
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, SW, SH, fill_color=AWS_DARK)
title_accent_bar(slide)

section_label(slide, "LIVE DEMO — FEATURE DETAIL")
add_textbox(slide, Inches(0.55), Inches(0.44), Inches(9), Inches(0.48),
            "Feature Detail Page — Deep-Dive Into Any Announcement", 22, bold=True, color=WHITE)

slide.shapes.add_picture(f"{ASSETS}/screen_detail.png",
                          Inches(0.35), Inches(1.0), Inches(12.65), Inches(5.95))

callouts2 = [
    (Inches(0.4),  Inches(1.05), "Context badges: date · service · domain · type"),
    (Inches(0.4),  Inches(1.55), "3–5 sentence AI summary, engineer-focused"),
    (Inches(0.4),  Inches(2.1),  "Structured key points for rapid evaluation"),
    (Inches(9.1),  Inches(1.05), "Link to original AWS announcement"),
    (Inches(9.1),  Inches(1.55), "Back button preserves search state"),
    (Inches(9.1),  Inches(2.1),  "Summarized-at timestamp for auditability"),
]
for cx, cy, text in callouts2:
    add_rect(slide, cx, cy, Inches(3.1), Inches(0.3), fill_color=RGBColor(0x23,0x2F,0x3E))
    add_textbox(slide, cx+Inches(0.08), cy+Inches(0.04), Inches(3.0), Inches(0.26),
                "▶  " + text, 9, color=AWS_ORANGE)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Screenshot: Filtered / Search Results
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, SW, SH, fill_color=AWS_DARK)
title_accent_bar(slide)

section_label(slide, "LIVE DEMO — FILTER & SEARCH IN ACTION")
add_textbox(slide, Inches(0.55), Inches(0.44), Inches(9), Inches(0.48),
            "Filtered View — 'bedrock' · service: amazon-bedrock · type: new-feature", 22, bold=True, color=WHITE)

slide.shapes.add_picture(f"{ASSETS}/screen_filtered.png",
                          Inches(0.35), Inches(1.0), Inches(12.65), Inches(5.95))

callouts3 = [
    (Inches(0.4),  Inches(1.05), "Keyword search highlights matching records"),
    (Inches(0.4),  Inches(1.55), "Active filter pills shown in header"),
    (Inches(0.4),  Inches(2.1),  "Result count updates in real-time"),
    (Inches(9.1),  Inches(1.05), "URL captures all filters — fully shareable"),
    (Inches(9.1),  Inches(1.55), "Infinite scroll loads next page automatically"),
    (Inches(9.1),  Inches(2.1),  "Works across all 8 feature type categories"),
]
for cx, cy, text in callouts3:
    add_rect(slide, cx, cy, Inches(3.1), Inches(0.3), fill_color=RGBColor(0x23,0x2F,0x3E))
    add_textbox(slide, cx+Inches(0.08), cy+Inches(0.04), Inches(3.0), Inches(0.26),
                "▶  " + text, 9, color=AWS_ORANGE)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Admin Use Cases
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, SW, SH, fill_color=LIGHT_GRAY)
add_rect(slide, 0, 0, SW, Inches(1.1), fill_color=AWS_DARK)
title_accent_bar(slide)

add_textbox(slide, Inches(0.55), Inches(0.2), Inches(11), Inches(0.56),
            "Real-World Use Cases for AWS Admins", 27, bold=True, color=WHITE)
add_textbox(slide, Inches(0.55), Inches(0.76), Inches(11), Inches(0.28),
            "From reactive monitoring to proactive innovation — the tracker enables both.", 13,
            color=AWS_ORANGE, italic=True)

use_cases = [
    (AWS_ORANGE, "Weekly Innovation Review",
     "Every Monday, the team opens the tracker filtered to 'last 7 days' to review everything "
     "AWS shipped. In 15 minutes instead of 2 hours, the team identifies 3-5 features worth piloting.",
     "Filter: from_date=7 days ago · sort=newest"),
    (AWS_BLUE, "Pre-Migration Research",
     "Before migrating workloads to Graviton4, search 'graviton4' to see all enhancements, "
     "GA announcements, and any related price changes — in one scrollable view.",
     "Search: 'graviton4' · type: general-availability"),
    (GREEN, "Deprecation Tracking",
     "Filter feature_type=deprecation to build a living list of EOL services and features. "
     "Share the URL with the architecture team as a standing reference.",
     "Filter: type=deprecation · sort=oldest"),
    (ACCENT_TEAL, "Cost Optimization Alerts",
     "Filter feature_type=price-change to immediately see DynamoDB, Lambda, and S3 price "
     "reductions. Calculate savings and present the business case for migration.",
     "Filter: type=price-change"),
    (RGBColor(0x8B,0x5C,0xF6), "New Service Evaluation",
     "When leadership asks 'should we adopt VPC Lattice?', search the tracker to see every "
     "feature shipped in the past year, understand the maturity trajectory.",
     "Search: 'vpc lattice' · date range: 2025-2026"),
    (RGBColor(0xEF,0x44,0x44), "Region Expansion Planning",
     "Filter feature_type=region-expansion for the services you use. Know the moment a "
     "service you rely on reaches your target region before checking the console.",
     "Filter: type=region-expansion · service=amazon-eks"),
]

COL_W = Inches(4.1)
ROW_H = Inches(2.0)
for i, (color, title, body, query) in enumerate(use_cases):
    col = i % 3
    row = i // 3
    bx = Inches(0.45 + col * (COL_W + Inches(0.12)))
    by = Inches(1.3 + row * (ROW_H + Inches(0.08)))

    add_rect(slide, bx, by, COL_W, ROW_H, fill_color=WHITE,
             line_color=RGBColor(0xE2,0xE8,0xF0), line_width_pt=1)
    add_rect(slide, bx, by, COL_W, Inches(0.08), fill_color=color)

    add_textbox(slide, bx+Inches(0.12), by+Inches(0.14), COL_W-Inches(0.2), Inches(0.3),
                title, 13, bold=True, color=DARK_GRAY)
    add_textbox(slide, bx+Inches(0.12), by+Inches(0.48), COL_W-Inches(0.2), Inches(0.88),
                body, 9.5, color=RGBColor(0x64,0x74,0x8B), wrap=True)

    # query pill
    add_rect(slide, bx+Inches(0.12), by+ROW_H-Inches(0.38), COL_W-Inches(0.24), Inches(0.28),
             fill_color=RGBColor(0xF1,0xF5,0xF9))
    add_textbox(slide, bx+Inches(0.18), by+ROW_H-Inches(0.36), COL_W-Inches(0.3), Inches(0.26),
                query, 8.5, color=color, italic=True)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — AI Summary Deep-Dive
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, SW, SH, fill_color=AWS_DARK)
title_accent_bar(slide)

section_label(slide, "AI INTELLIGENCE LAYER")
add_textbox(slide, Inches(0.55), Inches(0.44), Inches(10), Inches(0.52),
            "How Claude Haiku Transforms Raw AWS Announcements", 25, bold=True, color=WHITE)

# Left: Raw input
add_rect(slide, Inches(0.45), Inches(1.1), Inches(5.6), Inches(5.8),
         fill_color=RGBColor(0x1A,0x2A,0x3A), line_color=RGBColor(0x2D,0x3F,0x55), line_width_pt=1)
add_textbox(slide, Inches(0.6), Inches(1.15), Inches(5.3), Inches(0.3),
            "RAW RSS INPUT", 9, bold=True, color=MID_GRAY)
raw = (
    "<title>Amazon Bedrock Knowledge Bases now supports semantic reranking for improved "
    "retrieval-augmented generation (RAG) accuracy</title>\n\n"
    "<description><![CDATA[<p>Amazon Bedrock Knowledge Bases now supports semantic "
    "reranking, a technique that improves the accuracy of retrieved results for "
    "Retrieval Augmented Generation (RAG) applications...</p><p>With semantic "
    "reranking, you can now configure Knowledge Bases to use a cross-encoder model "
    "to rerank the initially retrieved documents or chunks based on their relevance "
    "to the query. This additional reranking step helps to improve the relevancy of "
    "retrieved results, which can reduce hallucinations and improve the quality of "
    "responses generated by foundation models...</p>]]></description>\n\n"
    "<category>general:products/amazon-bedrock</category>\n"
    "<category>marketing:marchitecture/ai-ml</category>"
)
add_textbox(slide, Inches(0.6), Inches(1.52), Inches(5.3), Inches(4.9),
            raw, 8.5, color=RGBColor(0x7D,0xD3,0xFC), wrap=True)

# Arrow
add_textbox(slide, Inches(6.2), Inches(3.6), Inches(1.0), Inches(0.5),
            "→", 36, bold=True, color=AWS_ORANGE, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(6.0), Inches(4.1), Inches(1.35), Inches(0.32),
            "Claude\nHaiku", 9, color=AWS_ORANGE, align=PP_ALIGN.CENTER)

# Right: AI output
add_rect(slide, Inches(7.3), Inches(1.1), Inches(5.6), Inches(5.8),
         fill_color=RGBColor(0x0D,0x2A,0x1A), line_color=GREEN, line_width_pt=1)
add_textbox(slide, Inches(7.5), Inches(1.15), Inches(5.3), Inches(0.3),
            "AI-ENRICHED OUTPUT", 9, bold=True, color=MID_GRAY)

output_sections = [
    ("feature_type:", "new-feature", AWS_ORANGE),
    ("summary:", (
        "Amazon Bedrock Knowledge Bases introduces semantic reranking to improve retrieval "
        "accuracy in RAG workflows. A cross-encoder model rescores retrieved chunks by relevance, "
        "significantly reducing hallucinations. Enable with a single API parameter. Supports "
        "Amazon Rerank 1.0 and Cohere Rerank 3.5."
    ), RGBColor(0xBB,0xF7,0xD0)),
    ("key_points:", (
        "• Reduces hallucinations by rescoring chunks\n"
        "• Supports Rerank 1.0 + Cohere Rerank 3.5\n"
        "• Single RetrievalConfiguration API param"
    ), RGBColor(0x86,0xEF,0xAC)),
]

ty = Inches(1.52)
for label, val, col in output_sections:
    add_textbox(slide, Inches(7.5), ty, Inches(5.3), Inches(0.26),
                label, 9, bold=True, color=MID_GRAY)
    ty += Inches(0.26)
    add_textbox(slide, Inches(7.5), ty, Inches(5.3), Inches(1.0),
                val, 10, color=col, wrap=True)
    ty += Inches(0.9)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Cost & Security
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, SW, SH, fill_color=LIGHT_GRAY)
add_rect(slide, 0, 0, SW, Inches(1.1), fill_color=AWS_DARK)
title_accent_bar(slide)

add_textbox(slide, Inches(0.55), Inches(0.2), Inches(11), Inches(0.56),
            "Cost-Efficient, Secure, and Production-Ready", 27, bold=True, color=WHITE)
add_textbox(slide, Inches(0.55), Inches(0.76), Inches(11), Inches(0.28),
            "Minimal cost, FedRAMP-aligned services, least-privilege IAM, and full observability.", 13,
            color=AWS_ORANGE, italic=True)

# Left: cost table
add_rect(slide, Inches(0.45), Inches(1.2), Inches(5.8), Inches(4.8),
         fill_color=WHITE, line_color=RGBColor(0xE2,0xE8,0xF0), line_width_pt=1)
add_textbox(slide, Inches(0.6), Inches(1.3), Inches(5.5), Inches(0.3),
            "Monthly Cost Estimate (~500 features/month)", 12, bold=True, color=DARK_GRAY)

rows = [
    ("Service", "Monthly Cost", True),
    ("Lambda (3 functions, arm64)", "~$0.01", False),
    ("DynamoDB on-demand", "~$0.02", False),
    ("S3 (summaries + SPA assets)", "~$0.03", False),
    ("CloudFront (PriceClass_100)", "~$0.10", False),
    ("Bedrock Claude Haiku (500 summaries)", "~$0.15", False),
    ("SQS + EventBridge Scheduler", "$0.00 (free tier)", False),
    ("Total", "~$0.31/month", True),
]
for i, (svc, cost, bold) in enumerate(rows):
    ry = Inches(1.72 + i * 0.44)
    if i % 2 == 0 and not bold:
        add_rect(slide, Inches(0.45), ry, Inches(5.8), Inches(0.44),
                 fill_color=RGBColor(0xF8,0xFA,0xFC))
    if bold:
        add_rect(slide, Inches(0.45), ry, Inches(5.8), Inches(0.44),
                 fill_color=RGBColor(0xFF,0xF7,0xED))
    add_textbox(slide, Inches(0.6), ry+Inches(0.1), Inches(3.8), Inches(0.28),
                svc, 10, bold=bold, color=DARK_GRAY)
    add_textbox(slide, Inches(4.5), ry+Inches(0.1), Inches(1.5), Inches(0.28),
                cost, 10, bold=bold,
                color=AWS_ORANGE if bold else GREEN, align=PP_ALIGN.RIGHT)

# Right: security & compliance
add_rect(slide, Inches(6.6), Inches(1.2), Inches(6.25), Inches(2.2),
         fill_color=WHITE, line_color=RGBColor(0xE2,0xE8,0xF0), line_width_pt=1)
add_textbox(slide, Inches(6.75), Inches(1.3), Inches(6.0), Inches(0.3),
            "Security & Compliance", 13, bold=True, color=DARK_GRAY)

sec_items = [
    "All services FedRAMP Moderate authorized (us-east-1 commercial)",
    "IAM least-privilege: each Lambda has only the permissions it needs",
    "S3 Block Public Access ON — CloudFront OAC for SPA, no direct access",
    "No VPC required — all traffic via AWS service endpoints",
    "API Gateway public read-only (no write surface exposed)",
]
for j, item in enumerate(sec_items):
    add_icon_bullet(slide, Inches(6.75), Inches(1.68 + j*0.3), Inches(5.9), Inches(0.28),
                    "•", item, 9.5, icon_color=GREEN, text_color=DARK_GRAY)

# Right bottom: reliability
add_rect(slide, Inches(6.6), Inches(3.55), Inches(6.25), Inches(2.45),
         fill_color=WHITE, line_color=RGBColor(0xE2,0xE8,0xF0), line_width_pt=1)
add_textbox(slide, Inches(6.75), Inches(3.65), Inches(6.0), Inches(0.3),
            "Reliability & Observability", 13, bold=True, color=DARK_GRAY)

rel_items = [
    "SQS DLQ with 7-day retention — no announcement ever lost",
    "Up to 3 automatic retries on Bedrock failures before DLQ",
    "CloudWatch Alarm: Summarizer errors > 10/hr → alert",
    "CloudWatch Alarm: DLQ messages > 0 → alert",
    "CloudFront + Lambda arm64 = low-latency, globally available",
]
for j, item in enumerate(rel_items):
    add_icon_bullet(slide, Inches(6.75), Inches(4.06 + j*0.3), Inches(5.9), Inches(0.28),
                    "•", item, 9.5, icon_color=AWS_BLUE, text_color=DARK_GRAY)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Future Roadmap
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, SW, SH, fill_color=AWS_DARK)
title_accent_bar(slide)

section_label(slide, "WHAT'S NEXT")
add_textbox(slide, Inches(0.55), Inches(0.44), Inches(10), Inches(0.52),
            "Roadmap — Where We Can Take This", 25, bold=True, color=WHITE)

roadmap = [
    ("Q2 2026", AWS_ORANGE, [
        "SNS / Slack notifications — ping #aws-updates when new features match watchlist",
        "Service watchlist per admin — personalized feed based on your AWS footprint",
        "Weekly digest email: top 5 features relevant to your services",
    ]),
    ("Q3 2026", AWS_BLUE, [
        "Impact scoring — Claude estimates cost/effort/risk of adopting each feature",
        "Multi-account tagging — link features to the accounts/OUs they affect",
        "Internal comments & status tracking (evaluating / piloting / adopted / declined)",
    ]),
    ("Q4 2026", GREEN, [
        "Cost savings calculator — estimate savings from price-change announcements",
        "Deprecation calendar — automatic JIRA tickets for EOL features",
        "Executive dashboard — monthly report of adopted innovations and estimated savings",
    ]),
]

for i, (quarter, color, items) in enumerate(roadmap):
    bx = Inches(0.45 + i * 4.28)
    by = Inches(1.1)
    bw = Inches(4.05)
    bh = Inches(5.5)

    add_rect(slide, bx, by, bw, bh,
             fill_color=RGBColor(0x1A,0x26,0x35), line_color=color, line_width_pt=2)
    add_rect(slide, bx, by, bw, Inches(0.44), fill_color=color)
    add_textbox(slide, bx, by+Inches(0.06), bw, Inches(0.34),
                quarter, 16, bold=True, color=AWS_DARK, align=PP_ALIGN.CENTER)

    for j, item in enumerate(items):
        iy = by + Inches(0.6 + j*1.52)
        add_rect(slide, bx+Inches(0.22), iy, Inches(3.62), Inches(1.4),
                 fill_color=RGBColor(0x23,0x32,0x44))
        add_icon_bullet(slide, bx+Inches(0.32), iy+Inches(0.1), Inches(3.42), Inches(1.2),
                        "•", item, 11, icon_color=color,
                        text_color=RGBColor(0xCB,0xD5,0xE1))

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Call to Action / Closing
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, SW, SH, fill_color=AWS_DARK)
add_rect(slide, 0, 0, Inches(0.18), SH, fill_color=AWS_ORANGE)

# Background accent lines
for i in range(1, 9):
    add_rect(slide, Inches(i*13.33/8), 0, Pt(0.5), SH,
             fill_color=RGBColor(0x2D, 0x3F, 0x55))

add_textbox(slide, Inches(0.65), Inches(1.1), Inches(11.5), Inches(0.4),
            "The tracker is live today.", 24,
            bold=False, color=MID_GRAY, italic=True)

add_textbox(slide, Inches(0.65), Inches(1.65), Inches(11.5), Inches(1.0),
            "Start using it this week.", 52, bold=True, color=WHITE)

add_textbox(slide, Inches(0.65), Inches(2.82), Inches(11), Inches(0.42),
            "https://d1s77r5nnbdjdt.cloudfront.net", 20,
            color=AWS_ORANGE, bold=True)

# Three CTA boxes
ctas = [
    (AWS_ORANGE, "Try It Now", "Open the app, search for any service you manage, and find features you may have missed."),
    (AWS_BLUE,   "Add Your Watchlist", "Tell us the top 10 AWS services your team owns — we'll prioritize those in the next sprint."),
    (GREEN,      "Give Feedback", "What filters, integrations, or views would make this indispensable for your workflow?"),
]
for i, (color, title, body) in enumerate(ctas):
    bx = Inches(0.65 + i * 4.15)
    by = Inches(3.7)
    add_rect(slide, bx, by, Inches(3.9), Inches(2.4),
             fill_color=RGBColor(0x1A,0x26,0x35))
    add_rect(slide, bx, by, Inches(3.9), Inches(0.08), fill_color=color)
    add_textbox(slide, bx+Inches(0.18), by+Inches(0.18), Inches(3.6), Inches(0.32),
                title, 14, bold=True, color=WHITE)
    add_textbox(slide, bx+Inches(0.18), by+Inches(0.58), Inches(3.6), Inches(1.65),
                body, 10.5, color=RGBColor(0xCB,0xD5,0xE1), wrap=True)

# Bottom bar
add_rect(slide, 0, SH - Inches(0.5), SW, Inches(0.5), fill_color=RGBColor(0x14,0x21,0x30))
add_textbox(slide, 0, SH - Inches(0.46), SW, Inches(0.42),
            "AWS Features Tracker  ·  Built on AWS SAM · Lambda · Bedrock · DynamoDB · CloudFront  ·  ~$0.31/month  ·  April 2026",
            9, color=MID_GRAY, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
